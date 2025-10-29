 # src/extract_pptx.py
from pathlib import Path
from hashlib import sha1
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def _pt(val):
    try:
        return float(val.pt)
    except Exception:
        return None

def _rgb(font):
    try:
        if font.color and font.color.rgb:
            return str(font.color.rgb)
    except Exception:
        pass
    return None

def extract_features_pptx(path: Path) -> pd.DataFrame:
    prs = Presentation(path)
    rows = []

    for s_idx, slide in enumerate(prs.slides, start=1):
        # text runs
        for shp in slide.shapes:
            if not hasattr(shp, "text_frame") or not shp.has_text_frame:
                continue
            for p_idx, para in enumerate(shp.text_frame.paragraphs):
                p_font = para.font  # paragraph default
                for r_idx, run in enumerate(para.runs):
                    rf = run.font
                    fam = rf.name or p_font.name
                    size = _pt(rf.size) or _pt(p_font.size)
                    rows.append({
                        "file": path.name,
                        "kind": "pptx_run",
                        "slide_idx": s_idx,
                        "shape_id": getattr(shp, "shape_id", None),
                        "para_idx": p_idx,
                        "run_idx": r_idx,
                        "text": run.text,
                        "font_family": fam,
                        "font_size_pt": size,
                        "bold": bool(rf.bold) if rf.bold is not None else None,
                        "italic": bool(rf.italic) if rf.italic is not None else None,
                        "underline": bool(rf.underline) if rf.underline is not None else None,
                        "color_rgb": _rgb(rf) or _rgb(p_font),
                    })

        # images
        for shp in slide.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shp.image.blob
                    rows.append({
                        "file": path.name,
                        "kind": "pptx_image",
                        "slide_idx": s_idx,
                        "shape_id": getattr(shp, "shape_id", None),
                        "sha1": sha1(blob).hexdigest(),
                        "text": None,
                        "font_family": None,
                        "font_size_pt": None,
                        "bold": None,
                        "italic": None,
                        "underline": None,
                        "color_rgb": None,
                    })
                except Exception:
                    pass

    return pd.DataFrame(rows)

